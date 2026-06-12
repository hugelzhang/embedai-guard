
"""EmbedAI Guard PPT Generator - Utility Functions"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG_DARK      = RGBColor(0x0A, 0x0E, 0x17)
BG_CARD      = RGBColor(0x12, 0x18, 0x24)
ACCENT_CYAN  = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_BLUE  = RGBColor(0x00, 0x96, 0xFF)
ACCENT_PURPLE= RGBColor(0xA7, 0x4B, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x7A)
ACCENT_ORANGE= RGBColor(0xFF, 0x8C, 0x00)
TEXT_WHITE   = RGBColor(0xE8, 0xEC, 0xF1)
TEXT_GRAY    = RGBColor(0x8A, 0x94, 0xA6)
TEXT_DIM     = RGBColor(0x4A, 0x55, 0x66)
LINE_DARK    = RGBColor(0x1E, 0x28, 0x36)

def set_slide_bg(slide, color=BG_DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_CYAN, height=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def TB(slide, left, top, width, height, text, font_size=Pt(14),
       color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def MLB(slide, left, top, width, height, lines, font_size=Pt(12),
        color=TEXT_WHITE, line_spacing=Pt(20), font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = line_spacing
    return txBox

def metric_card(slide, left, top, width, height, number, label, color=ACCENT_CYAN):
    add_rect(slide, left, top, width, height)
    TB(slide, left + Pt(16), top + Pt(12), width - Pt(32), Pt(36), number, Pt(32), color, True)
    TB(slide, left + Pt(16), top + height - Pt(30), width - Pt(32), Pt(20), label, Pt(11), TEXT_GRAY)

def section_title(slide, title, subtitle='', accent_color=ACCENT_CYAN):
    add_accent_line(slide, Inches(0.8), Inches(0.55), Inches(0.6), accent_color)
    TB(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.45), title, Pt(28), TEXT_WHITE, True)
    if subtitle:
        TB(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.35), subtitle, Pt(13), TEXT_GRAY)

def footer(slide, page_num, section=''):
    TB(slide, Inches(0.8), Inches(7.0), Inches(6), Inches(0.3), 'EmbedAI Guard  |  ' + section, Pt(8), TEXT_DIM)
    TB(slide, Inches(11.8), Inches(7.0), Inches(1), Inches(0.3), page_num, Pt(8), TEXT_DIM, alignment=PP_ALIGN.RIGHT)

def phase_card(slide, left, top, width, height, phase_num, title, period, items, color=ACCENT_CYAN):
    add_rect(slide, left, top, width, height, border_color=color)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Pt(12), top + Pt(12), Pt(28), Pt(28))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(phase_num)
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BG_DARK
    p.alignment = PP_ALIGN.CENTER
    TB(slide, left + Pt(48), top + Pt(14), width - Pt(60), Pt(22), title, Pt(14), TEXT_WHITE, True)
    TB(slide, left + Pt(12), top + Pt(40), width - Pt(24), Pt(16), period, Pt(9), color)
    MLB(slide, left + Pt(14), top + Pt(58), width - Pt(28), height - Pt(70),
        ['▸ ' + item for item in items], Pt(9), TEXT_GRAY, Pt(14))

def comp_table(slide, left, top, col_widths, headers, rows, hc=TEXT_WHITE, rc=None):
    if rc is None:
        rc = [TEXT_GRAY] * len(rows)
    y = top
    x = left
    for h, w in zip(headers, col_widths):
        TB(slide, x, y, w, Pt(20), h, Pt(10), hc, True)
        x += w
    y += Pt(18)
    add_accent_line(slide, left, y, sum(col_widths), LINE_DARK, Pt(1))
    y += Pt(8)
    for row, rcolor in zip(rows, rc):
        x = left
        for cell, w in zip(row, col_widths):
            TB(slide, x, y, w, Pt(18), cell, Pt(10), rcolor)
            x += w
        y += Pt(16)

def tag_rect(slide, left, top, width, height, text, color=ACCENT_CYAN):
    shape = add_rect(slide, left, top, width, height, BG_CARD, color)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
