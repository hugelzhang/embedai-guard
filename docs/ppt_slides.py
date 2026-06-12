"""
EmbedAI Guard - Slide Deck Generator
生成 14 页科技风 PPT
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from ppt_utils import *
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════ S1: COVER ═══════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_accent_line(s, Inches(0.8), Inches(2.2), Inches(1.8), ACCENT_CYAN, Pt(4))
add_accent_line(s, Inches(0.8), Inches(5.8), Inches(11.7), LINE_DARK, Pt(1))
TB(s, Inches(0.8), Inches(2.5), Inches(12), Inches(0.7),
   'EmbedAI Guard', Pt(54), TEXT_WHITE, True)
TB(s, Inches(0.8), Inches(3.15), Inches(12), Inches(0.5),
   '嵌入式固件开发的 AI Native 质量基础设施', Pt(22), ACCENT_CYAN)
TB(s, Inches(0.8), Inches(4.0), Inches(8), Inches(0.5),
   '三层守卫引擎 — 从「人工逐行审查」到「AI 自主闭环验证」', Pt(14), TEXT_GRAY)
for i, t in enumerate(['约束扫描', '契约验证', 'Golden Trace', '自动修复', '芯片插件']):
    tag_rect(s, Inches(0.8 + i * 2.3), Inches(5.0), Inches(2.0), Inches(0.4), t, ACCENT_CYAN)
TB(s, Inches(0.8), Inches(6.0), Inches(6), Inches(0.3),
   '项目计划书 v0.1  |  2026.06', Pt(10), TEXT_DIM)
TB(s, Inches(0.8), Inches(6.3), Inches(6), Inches(0.3),
   '状态：规划阶段  |  定位：嵌入式 + AI 基础设施', Pt(10), TEXT_DIM)
print('Slide 1 done')

# ═══════════════════════════════════════════ S2: PROBLEM ═══════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '问题陈述', '嵌入式开发的核心矛盾', ACCENT_ORANGE)
pains = [
    ('[速度]', '软件速度 x100，硬件速度不变',
     ['AI 秒级生成代码', '烧录/测试/调试仍然分钟级', 'AI 代码越多 — 瓶颈越突出']),
    ('[效率]', '60% 时间花在 Debug',
     ['引脚冲突、时序错误、驱动 Bug', '大量问题是已知规则的违规', '理论上完全可以自动化']),
    ('[工具]', '现有工具只管「发现」，不管「修复」',
     ['PC-lint 告诉你违规', 'Saleae 给你看波形', '但没人帮你修代码']),
]
for i, (icon, title, items) in enumerate(pains):
    x = Inches(0.8 + i * 4.1)
    add_rect(s, x, Inches(1.6), Inches(3.8), Inches(3.2))
    TB(s, x + Pt(16), Inches(1.7), Inches(3.4), Inches(0.35), icon, Pt(28), ACCENT_ORANGE, True)
    TB(s, x + Pt(16), Inches(2.2), Inches(3.4), Inches(0.3), title, Pt(16), TEXT_WHITE, True)
    MLB(s, x + Pt(16), Inches(2.7), Inches(3.3), Inches(1.8),
        ['> ' + item for item in items], Pt(12), TEXT_GRAY, Pt(24))
add_accent_line(s, Inches(0.8), Inches(5.2), Inches(11.7), LINE_DARK, Pt(1))
TB(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.6),
   '嵌入式是 AI 工具化最落后的领域之一 — 不是因为不需要，是因为没人认真做。',
   Pt(18), ACCENT_ORANGE, True)
footer(s, '02', '问题')
print('Slide 2 done')

# ═══════════════════════════════════════════ S3: MARKET ═══════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '市场机会', '嵌入式开发者 400-600 万 | 工具支出 $200-2000/年')
metrics = [
    ('400-600 万', '全球嵌入式开发者', ACCENT_CYAN),
    ('30+ 家', '国产 MCU 原厂', ACCENT_BLUE),
    ('$8 亿/年', '可触达市场规模', ACCENT_PURPLE),
    ('0 个', 'AI 原生竞品', ACCENT_GREEN),
]
for i, (num, label, color) in enumerate(metrics):
    x = Inches(0.8 + i * 3.1)
    metric_card(s, x, Inches(1.6), Inches(2.8), Inches(1.4), num, label, color)

TB(s, Inches(0.8), Inches(3.4), Inches(5), Inches(0.3),
   '核心客户：MCU 原厂（最大的买单方）', Pt(16), TEXT_WHITE, True)
TB(s, Inches(0.8), Inches(3.8), Inches(5), Inches(0.3),
   '原厂 SDK 是巨大的成本中心，直接影响芯片销量', Pt(12), TEXT_GRAY)
cust_rows = [
    ['MCU 原厂', 'SDK 质量低 — FAE 成本高 — 芯片丢单', '30-80 万/年', '★★★★★'],
    ['方案公司', '项目周期紧 — Debug 占 60%', '2-5 万/年', '★★★★'],
    ['创业团队', '人少活多 — 固件是瓶颈', '0.5-2 万/年', '★★★★'],
    ['个人开发者', '工具贵 — 学习成本高', '99 元/月', '★★'],
]
cust_colors = [ACCENT_CYAN, TEXT_WHITE, TEXT_GRAY, TEXT_GRAY]
comp_table(s, Inches(0.8), Inches(4.3),
           [Inches(2.0), Inches(5.3), Inches(2.5), Inches(1.7)],
           ['客户', '痛点', '付费能力', '意愿'],
           cust_rows, TEXT_WHITE, cust_colors)
footer(s, '03', '市场')
print('Slide 3 done')

# ═══════════════════════════════════════════ S4: SOLUTION ═══════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '解决方案：三层守卫引擎', '不是「发现问题」，是「自动修复」', ACCENT_GREEN)
guards = [
    ('01', 'Constraint Scanner\n约束扫描引擎', ACCENT_CYAN,
     ['delay/malloc/乘除自动检测', 'ISR 复杂度分析', 'printf 条件编译验证', '— 自动修复违规代码']),
    ('02', 'Contract Validator\n契约验证引擎', ACCENT_BLUE,
     ['引脚分配冲突检测', '时钟树配置验证', 'DMA 通道映射检查', '— 烧录前拦截硬件错误']),
    ('03', 'Golden Trace\n波形回归引擎', ACCENT_PURPLE,
     ['逻辑分析仪自动采集波形', '黄金数据对比 + 差异定位', '时序偏差自动分析', '— 驱动变更的自动化回归']),
]
for i, (num, title, color, items) in enumerate(guards):
    x = Inches(0.8 + i * 4.1)
    TB(s, x, Inches(1.6), Inches(2), Inches(0.6),
       num, Pt(48), color, True, font_name='Consolas')
    TB(s, x, Inches(2.1), Inches(3.5), Inches(0.6),
       title, Pt(16), TEXT_WHITE, True)
    MLB(s, x, Inches(2.9), Inches(3.5), Inches(2),
        ['✓ ' + item for item in items], Pt(11), TEXT_GRAY, Pt(18))
add_accent_line(s, Inches(0.8), Inches(5.5), Inches(11.7), LINE_DARK, Pt(1))
TB(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
   '三大引擎协同：软件侧 100 倍提速 — 硬件侧最小化烧录 — 一次验证通过',
   Pt(16), ACCENT_GREEN, True)
footer(s, '04', '方案')
print('Slide 4 done')

# ═══════════════════════════════════════════ S5: PRODUCT ═══════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '产品形态', 'CLI + VS Code 扩展 + 硬件桥接器')
prod_cols = [
    ('CLI 核心引擎', 'Python 3.10+', ACCENT_CYAN,
     ['guard scan — 约束扫描', 'guard check — 契约验证', 'guard trace — 波形对比',
      'guard fix — 自动修复', 'guard test — Mock 生成']),
    ('VS Code 扩展', 'TypeScript + LSP', ACCENT_BLUE,
     ['实时违规高亮 (squiggly)', 'Quick Fix 一键修复', '引脚冲突可视化面板',
      '波形对比视图', '状态栏质量指示器']),
    ('硬件桥接层', 'Saleae + pyOCD', ACCENT_PURPLE,
     ['Saleae Logic 自动采集', 'pyOCD / OpenOCD 烧录', '串口日志实时监控',
      '芯片插件热加载', 'JUnit XML 结果输出']),
]
for i, (title, tech, color, items) in enumerate(prod_cols):
    x = Inches(0.8 + i * 4.1)
    add_rect(s, x, Inches(1.5), Inches(3.8), Inches(3.8), border_color=color)
    add_accent_line(s, x, Inches(1.5), Inches(3.8), color, Pt(3))
    TB(s, x + Pt(16), Inches(1.7), Inches(3.4), Inches(0.3),
       title, Pt(16), TEXT_WHITE, True)
    TB(s, x + Pt(16), Inches(2.0), Inches(3.4), Inches(0.2),
       tech, Pt(10), color)
    MLB(s, x + Pt(16), Inches(2.5), Inches(3.3), Inches(2.5),
        ['▸ ' + item for item in items], Pt(11), TEXT_GRAY, Pt(20))
footer(s, '05', '产品')
print('Slide 5 done')

# ═══════════════════════════════════════════ S6: ARCHITECTURE ═══════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '技术架构', '插件化 · 可扩展 · 芯片无关核心层')
arch = [
    (Inches(1.5), 'VS Code Extension Layer', ACCENT_CYAN,
     '实时违规高亮  |  冲突面板  |  波形视图  |  修复建议'),
    (Inches(2.7), 'CLI Core Engine  (Python 3.10+)', ACCENT_BLUE,
     'ConstraintEngine  |  ContractEngine  |  GoldenTraceEngine  |  AutoFix Engine'),
    (Inches(3.9), 'Hardware Bridge Layer', ACCENT_PURPLE,
     'Saleae SDK  |  pyOCD / OpenOCD  |  Serial Monitor'),
    (Inches(5.1), 'Chip Plugin Layer  (JSON Schema)', ACCENT_GREEN,
     'STM32L4  |  HC32F003  |  YS32F003  |  ... Custom'),
]
for y, title, color, items in arch:
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.9), border_color=color)
    add_accent_line(s, Inches(0.8), y, Inches(0.08), color, Inches(0.9))
    TB(s, Inches(1.2), y + Pt(8), Inches(5), Inches(0.3),
       title, Pt(13), TEXT_WHITE, True)
    TB(s, Inches(1.2), y + Pt(30), Inches(10), Inches(0.25),
       items, Pt(10), TEXT_GRAY)
TB(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.3),
   '核心技术：tree-sitter C AST  ·  YAML 规则 DSL  ·  JSON Schema 芯片契约  ·  Saleae Python SDK  ·  VS Code LSP',
   Pt(10), TEXT_DIM)
footer(s, '06', '架构')
print('Slide 6 done')

# ═══════════════════════════════════════════ S7: RULES DSL ═══════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '约束规则 DSL', '编码规范 — 可执行规则 — 自动检测 + 自动修复')
rules_rows = [
    ['EMBED-001', '禁止阻塞延时', 'error', 'HAL_Delay 等阻塞函数', '定时器 + 状态机'],
    ['EMBED-002', '禁止动态内存', 'error', 'malloc / free 调用', '静态数组 / 对象池'],
    ['EMBED-003', '禁止乘除', 'warning', 'AST binary_expression /*', '/2 — >>1 (自动)'],
    ['EMBED-004', 'ISR 简洁性', 'error', 'IRQHandler 行数 > 10', '拆分 ISR + 主循环'],
    ['EMBED-005', '禁止浮点', 'warning', 'float / double 声明', '定点整数运算 (mV单位)'],
    ['EMBED-006', 'printf 条件编译', 'warning', 'printf 未包裹 DEBUG', '替换为项目 LOG 宏'],
]
comp_table(s, Inches(0.8), Inches(1.6),
           [Inches(1.3), Inches(1.8), Inches(1.0), Inches(3.0), Inches(3.5)],
           ['ID', '规则', '级别', '检测方式', '修复策略'],
           rules_rows)
add_accent_line(s, Inches(0.8), Inches(4.3), Inches(11.7), LINE_DARK, Pt(1))
TB(s, Inches(0.8), Inches(4.5), Inches(5), Inches(0.25),
   '规则 DSL 示例 (YAML)', Pt(14), TEXT_WHITE, True)
code_lines = [
    'id: EMBED-001',
    'name: 禁止阻塞延时',
    'severity: error',
    'patterns:',
    '  - HAL_Delay(',
    '  - delay_ms(',
    'fix_strategy:',
    '  type: suggest',
    '  template:',
    '    改为定时器驱动：',
    '    static uint32_t tick = 0;',
    '    if (HAL_GetTick() - tick >= ms)',
]
MLB(s, Inches(0.8), Inches(4.85), Inches(5.5), Inches(1.8),
    code_lines, Pt(9), ACCENT_CYAN, Pt(4), font_name='Consolas')
TB(s, Inches(7.0), Inches(4.5), Inches(5), Inches(0.25),
   '检测 + 修复流程', Pt(14), TEXT_WHITE, True)
flow_lines = [
    '1. tree-sitter 解析 C 源码 -> AST',
    '2. 遍历 AST 节点，匹配规则 pattern',
    '3. 匹配成功 -> 报告违规位置和级别',
    '4. fix_strategy.type == auto ?',
    '   -> 自动变换代码 -> 重新扫描',
    '5. 输出：终端报告 + JSON + JUnit XML',
]
MLB(s, Inches(7.0), Inches(4.85), Inches(5.5), Inches(1.8),
    flow_lines, Pt(10), TEXT_GRAY, Pt(16))
footer(s, '07', '规则')
print('Slide 7 done')

# ═══════════════════════════════════════════ S8: CONTRACT ═══════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '芯片契约系统',
              'JSON Schema — 引脚 / 时钟 / DMA / 中断 四维冲突检测')
quads = [
    ('引脚冲突检测', ACCENT_CYAN, [
        '同一引脚分配两个外设 — error',
        'ADC 通道重复映射 — error',
        'SWD 调试口被复用 — warning',
        '5V 容忍引脚接 5V 信号 — info']),
    ('时钟树验证', ACCENT_BLUE, [
        'PLL 倍频超出 VCO 范围 — error',
        'APB1/2 时钟超限 — error',
        'HSE 频率不在 [4,48]MHz — error',
        'Flash 等待周期不匹配 — warning']),
    ('DMA 映射检查', ACCENT_PURPLE, [
        '两外设争用同一 DMA 通道 — error',
        'DMA 请求映射不存在 — error',
        '方向与数据宽度不匹配 — warning',
        'FIFO 阈值配置错误 — warning']),
    ('中断优先级验证', ACCENT_GREEN, [
        'FreeRTOS ISR 优先级溢出 — error',
        '优先级分组不一致 — warning',
        '同级优先级无抢占 — info',
        '中断嵌套层级过深 — warning']),
]
for i, (title, color, items) in enumerate(quads):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.1)
    y = Inches(1.5 + row * 2.6)
    add_rect(s, x, y, Inches(5.8), Inches(2.3), border_color=color)
    add_accent_line(s, x, y, Inches(5.8), color, Pt(3))
    TB(s, x + Pt(14), y + Pt(10), Inches(5.4), Inches(0.25),
       title, Pt(15), TEXT_WHITE, True)
    MLB(s, x + Pt(14), y + Pt(40), Inches(5.4), Inches(1.8),
        items, Pt(10), TEXT_GRAY, Pt(16))
footer(s, '08', '契约')
print('Slide 8 done')

# ═══════════════════════════════════════════ S9: COMPETITION ═══════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '竞争格局',
              '现有工具全部是「发现问题」，没有一个能「自动修复」')
comp_rows = [
    ['EmbedAI Guard', '✓', '✓', '✓', '✓', '免费-定制'],
    ['PC-lint', '✓', '✗', '✗', '✗', '$399'],
    ['MISRA-C Checker', '✓', '✗', '✗', '✗', '$$$$'],
    ['SonarQube', '✓', '✗', '✗', '✗', '$150+/年'],
    ['Copilot/Cursor', '~部分', '~部分', '✗', '✗', '$10-20/月'],
    ['Zephyr Twister', '✓ CI', '✗', '~仿真', '✓', '开源'],
    ['Renode', '✗', '✗', '✓ 仿真', '✓', '开源'],
    ['Saleae Logic', '✗', '✗', '✓ 采集', '✗', '$499-1499'],
]
comp_colors = [ACCENT_CYAN] + [TEXT_GRAY]*7
comp_table(s, Inches(0.8), Inches(1.6),
           [Inches(2.3), Inches(1.2), Inches(1.3), Inches(1.3), Inches(1.5), Inches(1.8)],
           ['产品', '检测', '修复', '硬件验证', '芯片识别', '价格'],
           comp_rows, TEXT_WHITE, comp_colors)
add_accent_line(s, Inches(0.8), Inches(4.6), Inches(11.7), LINE_DARK, Pt(1))
TB(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.4),
   '核心差异化：唯一同时具备「自动修复 + 硬件验证 + 芯片感知」能力的产品',
   Pt(18), ACCENT_GREEN, True)
MLB(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.8), [
    '定位：不是替代 Copilot — 是在 Copilot 输出上增加一层嵌入式专用质量门禁',
    '类比：Copilot = 写草稿的人  |  EmbedAI Guard = 审核草稿的资深嵌入式工程师',
], Pt(13), TEXT_GRAY, Pt(20))
footer(s, '09', '竞争')
print('Slide 9 done')

# ═══════════════════════════════════════════ S10: BUSINESS MODEL ══════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '商业模式', '开源漏斗 — 个人转化 — 团队增长 — 企业签约')
pricing = [
    ('Community', '免费', ACCENT_GREEN,
     ['约束扫描 (10规则)', 'VS Code 扩展', '公开芯片插件', 'GitHub Issues']),
    ('Pro', '99 元/月', ACCENT_CYAN,
     ['全部规则 (50+)', '自动修复', 'Golden Trace 基础版', '优先邮件支持']),
    ('Team', '1999 元/月', ACCENT_BLUE,
     ['团队共享规则库', '私有芯片插件', 'Mock 测试生成', 'CI/CD 集成模板']),
    ('Enterprise', '定制报价', ACCENT_PURPLE,
     ['原厂 SDK 质量平台', '深度集成定制', '专属芯片插件', 'SLA + 专属 FAE']),
]
for i, (name, price, color, items) in enumerate(pricing):
    x = Inches(0.8 + i * 3.1)
    add_rect(s, x, Inches(1.5), Inches(2.8), Inches(3.0), border_color=color)
    add_accent_line(s, x, Inches(1.5), Inches(2.8), color, Pt(4))
    TB(s, x + Pt(14), Inches(1.7), Inches(2.4), Inches(0.25),
       name, Pt(18), TEXT_WHITE, True)
    TB(s, x + Pt(14), Inches(2.0), Inches(2.4), Inches(0.3),
       price, Pt(22), color, True)
    MLB(s, x + Pt(14), Inches(2.5), Inches(2.4), Inches(1.6),
        ['✓ ' + item for item in items], Pt(10), TEXT_GRAY, Pt(16))
TB(s, Inches(0.8), Inches(4.9), Inches(5), Inches(0.25),
   '3 年收入预估', Pt(14), TEXT_WHITE, True)
rev_data = [
    ['Y1', '200 Pro + 10 Team + 2 Enterprise', '148 万'],
    ['Y2', '500 Pro + 30 Team + 5 Enterprise', '432 万'],
    ['Y3', '1000 Pro + 60 Team + 10 Enterprise', '964 万'],
]
rev_colors = [ACCENT_CYAN, ACCENT_BLUE, ACCENT_PURPLE]
comp_table(s, Inches(0.8), Inches(5.3),
           [Inches(0.8), Inches(5.5), Inches(2.0)],
           ['年份', '客户构成', '年收入'],
           rev_data, TEXT_WHITE, rev_colors)
footer(s, '10', '商业')
print('Slide 10 done')

# ═══════════════════════════════════════════ S11: ROADMAP 1-3 ═══════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '开发路线图 Phase 1-3',
              '从 MVP 到商业化 · 9 个月 · 3 个可商用里程碑')
p1_3 = [
    ('P1', '约束引擎 MVP', '第 1-2 个月', ACCENT_GREEN, [
        'tree-sitter C AST 解析', '5 条核心规则 + 自动修复',
        'VS Code 扩展发布', 'GitHub 开源 + 技术博客']),
    ('P2', '契约引擎', '第 3-5 个月', ACCENT_CYAN, [
        '引脚/时钟/DMA/中断冲突检测', 'STM32L4/F1 芯片插件',
        '原厂 PoC 洽谈', '首份付费合同']),
    ('P3', 'Golden Trace', '第 6-9 个月', ACCENT_BLUE, [
        'Saleae 波形采集 + 对比', 'Mock 测试自动生成',
        '团队版发布', 'ARR > 30 万']),
]
for i, (p, t, period, color, items) in enumerate(p1_3):
    x = Inches(0.8 + i * 4.1)
    phase_card(s, x, Inches(1.5), Inches(3.8), Inches(4.2),
               p, t, period, items, color)
add_accent_line(s, Inches(0.8), Inches(5.9), Inches(11.7), LINE_DARK, Pt(1))
TB(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4),
   'P1+P2+P3 可独立商用  |  P4+P5 面向未来：规格即代码 — 数字孪生',
   Pt(14), ACCENT_CYAN, True)
footer(s, '11', '路线图')
print('Slide 11 done')

# ═══════════════════════════════════════════ S12: ROADMAP 4-5 ═══════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '中长期：规格即代码 — 数字孪生',
              '从「AI 辅助验证」到「AI 自主推导」')
p4_5 = [
    ('P4', '规格即代码', '第 10-15 个月', ACCENT_PURPLE, [
        'Datasheet 参数结构化提取', '寄存器映射自动生成',
        '时序约束自动验证', 'Renode/QEMU 仿真集成']),
    ('P5', '数字孪生', '第 16-24 个月', ACCENT_ORANGE, [
        '完整外设仿真模型', '多芯片协同仿真',
        '90% 验证在仿真中完成', '10+ 芯片插件 + 5+ 原厂客户']),
]
for i, (p, t, period, color, items) in enumerate(p4_5):
    x = Inches(0.8 + i * 6.1)
    phase_card(s, x, Inches(1.5), Inches(5.8), Inches(3.5),
               p, t, period, items, color)
add_accent_line(s, Inches(0.8), Inches(5.4), Inches(11.7), LINE_DARK, Pt(1))
TB(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
   '最终愿景', Pt(18), TEXT_WHITE, True)
MLB(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8), [
    '用户声明需求 — AI 读手册推导实现 — 数字孪生验证 — 一次性烧录 — 硬件闭环确认',
    '不是替代嵌入式工程师，是让工程师从重复劳动中解放，专注真正的创新',
], Pt(13), TEXT_GRAY, Pt(20))
footer(s, '12', '路线图')
print('Slide 12 done')

# ═══════════════════════════════════════════ S13: FINANCE & RISK ═══════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '启动成本 & 风险', '轻启动 · 快验证 · 步步为营')
TB(s, Inches(0.8), Inches(1.5), Inches(5), Inches(0.3),
   '启动成本（第 1-6 个月）', Pt(16), TEXT_WHITE, True)
cost_data = [
    ['开发工具', '0 元', '现有 VS Code + Python + 开源库'],
    ['硬件设备', '5,000 元', '逻辑分析仪 + 各芯片开发板'],
    ['云服务', '500 元/月', 'GitHub + 文档托管'],
    ['市场推广', '10,000 元', '技术博客 + 社区推广'],
    ['合计', '约 18,000 元', '盈亏平衡：第 6-8 个月'],
]
cost_colors = [ACCENT_GREEN, ACCENT_CYAN, ACCENT_CYAN, ACCENT_BLUE, ACCENT_GREEN]
comp_table(s, Inches(0.8), Inches(2.0),
           [Inches(1.5), Inches(1.5), Inches(4.0)],
           ['项目', '预算', '说明'],
           cost_data, TEXT_WHITE, cost_colors)

TB(s, Inches(7.0), Inches(1.5), Inches(5), Inches(0.3),
   '关键风险 & 应对', Pt(16), TEXT_WHITE, True)
risks = [
    ('原厂决策周期长', 'B2D 先行，个人/团队收入维持现金流'),
    ('竞品大厂入场', '深扎国产 MCU，大厂不做的垂直领域'),
    ('芯片 Plugin 维护成本高', '社区贡献 + 原厂自行维护'),
    ('AI 技术路线不确定', '约束引擎不依赖 LLM，规则驱动'),
    ('用户习惯不改变', '从 lint 替代品切入，渐进升级'),
]
for i, (risk, resp) in enumerate(risks):
    y = Inches(2.0 + i * 0.7)
    add_rect(s, Inches(7.0), y, Inches(5.5), Inches(0.6),
             border_color=LINE_DARK)
    TB(s, Inches(7.2), y + Pt(5), Inches(2.2), Inches(0.2),
       '[!] ' + risk, Pt(10), TEXT_GRAY)
    TB(s, Inches(7.2), y + Pt(24), Inches(5.0), Inches(0.2),
       '-> ' + resp, Pt(9), ACCENT_CYAN)
footer(s, '13', '财务/风险')
print('Slide 13 done')

# ═══════════════════════════════════════════ S14: NEXT STEPS ═══════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
section_title(s, '下一步行动', '本周可启动，2 个月出 MVP', ACCENT_GREEN)
steps = [
    ('Week 1-2', '约束扫描器核心', ACCENT_CYAN, [
        'tree-sitter C 集成', '5 条核心规则实现',
        'CLI: guard scan + guard fix']),
    ('Week 3-4', '自动修复引擎', ACCENT_BLUE, [
        'fix_strategy 框架', '3 条自动修复规则',
        '修复后自动重新验证']),
    ('Week 5-6', 'VS Code 扩展', ACCENT_PURPLE, [
        'LSP 实时违规高亮', 'Quick Fix 一键修复',
        'Marketplace 发布准备']),
    ('Week 7-8', '开源发布', ACCENT_GREEN, [
        'GitHub 仓库 + 文档', '示例项目 + 技术博客',
        '社区推广 + Star 目标 500+']),
]
for i, (week, title, color, items) in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    add_rect(s, x, Inches(1.5), Inches(2.8), Inches(2.8), border_color=color)
    add_accent_line(s, x, Inches(1.5), Inches(2.8), color, Pt(3))
    TB(s, x + Pt(14), Inches(1.7), Inches(2.4), Inches(0.2),
       week, Pt(10), color)
    TB(s, x + Pt(14), Inches(1.9), Inches(2.4), Inches(0.3),
       title, Pt(15), TEXT_WHITE, True)
    MLB(s, x + Pt(14), Inches(2.3), Inches(2.4), Inches(1.5),
        ['> ' + item for item in items], Pt(10), TEXT_GRAY, Pt(16))
add_accent_line(s, Inches(0.8), Inches(4.8), Inches(11.7), ACCENT_CYAN, Pt(3))
TB(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
   'Phase 1 目标：用 STM32L475 Pandora 项目作为验证平台，证明产品价值',
   Pt(18), TEXT_WHITE, True)
TB(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.7),
   '准备启动了吗？', Pt(32), ACCENT_GREEN, True)
footer(s, '14', '行动')
print('Slide 14 done')

# ═══════════════════════════════════════════ SAVE ═══════════════════════
output = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      'EmbedAI_Guard_Pitch.pptx')
prs.save(output)
print(f'\n=== DONE ===')
print(f'File: {output}')
print(f'Slides: {len(prs.slides)}')
